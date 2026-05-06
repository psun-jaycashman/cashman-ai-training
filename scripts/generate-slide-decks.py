"""
Generate Cashman AI Training slide decks (.pptx) for every module + lesson video.

Style is matched to docs/video-scripts/slides/pptx/_reference (CashmanAILandL.pptx):
  - Cream background  #FDFAF7
  - Cashman green     #2B6E2B (accent bars, links)
  - Light green tint  #C8E9C8 (callout band)
  - Petrona Bold for headlines, Inter for body
  - 16:9, generous left margin, left-aligned content

Each deck is short (2-3 slides) and acts as background for the AI avatar, which
sits in the bottom-right corner. No content is placed inside the AVATAR_ZONE.

Run:
    python3 scripts/generate-slide-decks.py
Output:
    docs/video-scripts/slides/pptx/module-N/<deck-slug>.pptx
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# --------------------------------------------------------------------------- #
# Theme
# --------------------------------------------------------------------------- #

CREAM = RGBColor(0xFD, 0xFA, 0xF7)
GREEN = RGBColor(0x2B, 0x6E, 0x2B)
GREEN_TINT = RGBColor(0xC8, 0xE9, 0xC8)
HEADLINE = RGBColor(0x03, 0x03, 0x03)
BODY = RGBColor(0x27, 0x25, 0x25)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)

HEADLINE_FONT = "Petrona"
BODY_FONT = "Inter"

SLIDE_W = Inches(16)
SLIDE_H = Inches(9)

# Avatar zone — bottom-right corner. Nothing important goes here.
AVATAR_LEFT = Inches(11.5)
AVATAR_TOP = Inches(3.0)
AVATAR_W = Inches(4.5)
AVATAR_H = Inches(6.0)

# Content area (left of the avatar zone)
CONTENT_LEFT = Inches(0.87)
CONTENT_RIGHT = Inches(11.0)
CONTENT_W = Inches(10.13)


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #


def _add_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font=BODY_FONT,
    size_pt=18,
    color=BODY,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def _add_avatar_placeholder(slide, mark=False):
    """Reserve the avatar zone. By default we just leave it empty.
    If `mark=True`, draw a faint outline so the deck builder can see the zone."""
    if not mark:
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, AVATAR_LEFT, AVATAR_TOP, AVATAR_W, AVATAR_H
    )
    shape.fill.background()
    shape.line.color.rgb = MUTED
    shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    _add_text(
        slide,
        AVATAR_LEFT,
        AVATAR_TOP + Inches(2.7),
        AVATAR_W,
        Inches(0.4),
        "AVATAR ZONE",
        size_pt=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def _set_background(slide, color=CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# --------------------------------------------------------------------------- #
# Slide layouts
# --------------------------------------------------------------------------- #


def add_title_slide(
    prs: Presentation,
    eyebrow: str,
    headline: str,
    subhead: str,
):
    """Cover/intro slide. Big headline, subhead, small eyebrow tag."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_background(slide)

    # Left green accent bar
    _add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, GREEN)

    # Eyebrow tag (uppercase, small, green)
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(2.2),
        CONTENT_W,
        Inches(0.4),
        eyebrow.upper(),
        font=BODY_FONT,
        size_pt=14,
        color=GREEN,
        bold=True,
    )

    # Headline (Petrona Bold, large)
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(2.7),
        CONTENT_W,
        Inches(2.0),
        headline,
        font=HEADLINE_FONT,
        size_pt=54,
        color=HEADLINE,
        bold=True,
        line_spacing=1.05,
    )

    # Subhead
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(5.2),
        CONTENT_W,
        Inches(1.5),
        subhead,
        font=BODY_FONT,
        size_pt=22,
        color=BODY,
        line_spacing=1.25,
    )

    # Cashman wordmark in lower left
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(8.1),
        Inches(6.0),
        Inches(0.4),
        "CASHMAN  ·  AI TRAINING",
        font=BODY_FONT,
        size_pt=11,
        color=GREEN,
        bold=True,
    )

    _add_avatar_placeholder(slide)
    return slide


def add_concept_slide(
    prs: Presentation,
    eyebrow: str,
    headline: str,
    body_lines: Sequence[str],
):
    """Concept slide. Eyebrow + headline + bullet body. Used for terms,
    frameworks, definitions."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_background(slide)

    # Top eyebrow
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(0.9),
        CONTENT_W,
        Inches(0.4),
        eyebrow.upper(),
        font=BODY_FONT,
        size_pt=12,
        color=GREEN,
        bold=True,
    )

    # Headline
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(1.4),
        CONTENT_W,
        Inches(1.6),
        headline,
        font=HEADLINE_FONT,
        size_pt=42,
        color=HEADLINE,
        bold=True,
        line_spacing=1.1,
    )

    # Green divider rule
    _add_rect(slide, CONTENT_LEFT, Inches(3.2), Inches(0.8), Inches(0.06), GREEN)

    # Body bullets — bigger text, plenty of breathing room
    body_text = "\n".join(body_lines)
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(3.7),
        CONTENT_W,
        Inches(4.5),
        body_text,
        font=BODY_FONT,
        size_pt=22,
        color=BODY,
        line_spacing=1.45,
    )

    _add_avatar_placeholder(slide)
    return slide


def add_image_slide(
    prs: Presentation,
    eyebrow: str,
    headline: str,
    image_path: Path,
    caption: str = "",
):
    """Image-led slide. Eyebrow + short headline at top, real screenshot
    centered in the content area, optional caption underneath. The image
    is letterboxed into a fixed box so different screenshot aspect ratios
    still look intentional."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_background(slide)

    # Eyebrow
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(0.7),
        CONTENT_W,
        Inches(0.4),
        eyebrow.upper(),
        font=BODY_FONT,
        size_pt=12,
        color=GREEN,
        bold=True,
    )

    # Headline
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(1.15),
        CONTENT_W,
        Inches(1.1),
        headline,
        font=HEADLINE_FONT,
        size_pt=30,
        color=HEADLINE,
        bold=True,
        line_spacing=1.05,
    )

    # Image area (left of avatar zone). Fit-inside scale, top-anchored.
    box_left = CONTENT_LEFT
    box_top = Inches(2.45)
    box_w_in = 10.13
    box_h_in = 5.4

    from PIL import Image as PILImage  # local import; keeps top-of-file lean
    with PILImage.open(image_path) as im:
        img_w, img_h = im.size
    img_aspect = img_w / img_h
    box_aspect = box_w_in / box_h_in
    if img_aspect >= box_aspect:
        # image is wider than the box → fit to width
        draw_w_in = box_w_in
        draw_h_in = box_w_in / img_aspect
    else:
        draw_h_in = box_h_in
        draw_w_in = box_h_in * img_aspect
    # center horizontally, top-align vertically inside the box
    offset_x_in = (box_w_in - draw_w_in) / 2.0

    # Light card background behind the image so screenshots with dark
    # chrome (e.g. Cashman Portal) don't float against the cream slide.
    pad = Inches(0.12)
    _add_rect(
        slide,
        box_left + Inches(offset_x_in) - pad,
        box_top - pad,
        Inches(draw_w_in) + pad * 2,
        Inches(draw_h_in) + pad * 2,
        WHITE,
    )
    # Hairline green underline beneath the card for accent.
    _add_rect(
        slide,
        box_left + Inches(offset_x_in) - pad,
        box_top + Inches(draw_h_in) + pad,
        Inches(draw_w_in) + pad * 2,
        Inches(0.04),
        GREEN,
    )

    slide.shapes.add_picture(
        str(image_path),
        box_left + Inches(offset_x_in),
        box_top,
        Inches(draw_w_in),
        Inches(draw_h_in),
    )

    # Caption below the card
    if caption:
        _add_text(
            slide,
            CONTENT_LEFT,
            box_top + Inches(draw_h_in) + Inches(0.45),
            CONTENT_W,
            Inches(0.7),
            caption,
            font=BODY_FONT,
            size_pt=14,
            color=MUTED,
            italic=True,
            line_spacing=1.3,
        )

    _add_avatar_placeholder(slide)
    return slide


def add_callout_slide(
    prs: Presentation,
    eyebrow: str,
    headline: str,
    callout: str,
):
    """A bold takeaway slide with a green callout band. Used for rules,
    warnings, and the action hand-off at the end of a video."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_background(slide)

    # Eyebrow
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(0.9),
        CONTENT_W,
        Inches(0.4),
        eyebrow.upper(),
        font=BODY_FONT,
        size_pt=12,
        color=GREEN,
        bold=True,
    )

    # Headline
    _add_text(
        slide,
        CONTENT_LEFT,
        Inches(1.4),
        CONTENT_W,
        Inches(2.0),
        headline,
        font=HEADLINE_FONT,
        size_pt=42,
        color=HEADLINE,
        bold=True,
        line_spacing=1.1,
    )

    # Light green callout panel
    panel_top = Inches(4.2)
    panel_h = Inches(2.6)
    _add_rect(slide, CONTENT_LEFT, panel_top, CONTENT_W, panel_h, GREEN_TINT)
    _add_rect(
        slide,
        CONTENT_LEFT - Inches(0.04),
        panel_top,
        Inches(0.13),
        panel_h,
        GREEN,
    )

    _add_text(
        slide,
        CONTENT_LEFT + Inches(0.4),
        panel_top + Inches(0.35),
        CONTENT_W - Inches(0.8),
        panel_h - Inches(0.7),
        callout,
        font=HEADLINE_FONT,
        size_pt=26,
        color=HEADLINE,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.25,
    )

    _add_avatar_placeholder(slide)
    return slide


# --------------------------------------------------------------------------- #
# Deck definitions — content for every video
#
# Each deck = 2 or 3 slides. Slide types: title, concept, callout.
# These are SHORT — they exist as background while the avatar narrates.
# --------------------------------------------------------------------------- #


@dataclass
class Slide:
    kind: str  # "title" | "concept" | "callout" | "image"
    eyebrow: str = ""
    headline: str = ""
    subhead: str = ""
    body: list[str] = field(default_factory=list)
    callout: str = ""
    image: str = ""  # path relative to docs/video-scripts/slides/assets/
    caption: str = ""


@dataclass
class Deck:
    module: int
    slug: str
    title: str
    slides: list[Slide]


DECKS: list[Deck] = [
    # ===================================================================== #
    # MODULE 1 — Your AI Toolkit
    # ===================================================================== #
    Deck(
        module=1,
        slug="module-1-intro",
        title="Module 1 Intro — Your AI Toolkit",
        slides=[
            Slide(
                kind="title",
                eyebrow="Module 1",
                headline="Your AI Toolkit",
                subhead="Three tools. Three jobs. The simple rules that keep you safe.",
            ),
            Slide(
                kind="concept",
                eyebrow="The three tools",
                headline="Portal · Copilot · ChatGPT",
                body=[
                    "•  Cashman AI Portal — for company data",
                    "•  Microsoft Copilot — for Office work",
                    "•  ChatGPT — for public, general use",
                ],
            ),
        ],
    ),
    Deck(
        module=1,
        slug="lesson-1-1-what-is-ai",
        title="Lesson 1.1 — What is AI?",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 1",
                headline="What is AI?",
                subhead="A pattern-matching tool that learned to talk.",
            ),
            Slide(
                kind="concept",
                eyebrow="New term",
                headline="LLM — Large Language Model",
                body=[
                    "An AI that learned the patterns of language by reading billions of pages.",
                    "It predicts the next word — one at a time.",
                    "Same prompt twice = possibly different answers (it's non-deterministic).",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="The skill we're building",
                headline="Hallucinations",
                callout="AI confidently invents facts — dates, numbers, citations, names.\nAlways review the output.",
            ),
        ],
    ),
    Deck(
        module=1,
        slug="lesson-1-2-toolkit",
        title="Lesson 1.2 — Cashman AI Portal, Copilot, ChatGPT",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 2",
                headline="Your AI Toolkit",
                subhead="Three tools. Three jobs. One simple rule for each.",
            ),
            Slide(
                kind="concept",
                eyebrow="Which tool when",
                headline="Portal · Copilot · ChatGPT",
                body=[
                    "•  Portal  →  anything sensitive (contracts, bids, financials)",
                    "•  Copilot →  Outlook, Word, Excel, PowerPoint",
                    "•  ChatGPT →  general, public, non-sensitive learning & brainstorming",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="The rule",
                headline="The Bulletin Board Test",
                callout="If you wouldn't pin it on a public board in the lobby,\ndon't paste it into a public AI tool.",
            ),
        ],
    ),
    Deck(
        module=1,
        slug="lesson-1-3-ground-rules",
        title="Lesson 1.3 — The Ground Rules",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 3",
                headline="The Ground Rules",
                subhead="Four rules. Memorize them. Use them every day.",
            ),
            Slide(
                kind="concept",
                eyebrow="The four rules",
                headline="Verify. Protect. Own it.",
                body=[
                    "1.  Bulletin Board Test — sensitive data → Cashman AI Portal only.",
                    "2.  Always verify AI output — check facts, dates, numbers.",
                    "3.  You are responsible for what you send.",
                    "4.  AI augments. You decide.",
                ],
            ),
        ],
    ),
    # ===================================================================== #
    # MODULE 2 — AI and Email
    # ===================================================================== #
    Deck(
        module=2,
        slug="module-2-intro",
        title="Module 2 Intro — AI and Email",
        slides=[
            Slide(
                kind="title",
                eyebrow="Module 2",
                headline="AI and Email",
                subhead="Fix bad emails, summarize the chaos, write the hard ones.",
            ),
            Slide(
                kind="concept",
                eyebrow="Three exercises",
                headline="Email is the highest-ROI use of AI.",
                body=[
                    "1.  Fix a poorly written email",
                    "2.  Summarize a 47-reply project thread",
                    "3.  Draft a difficult client response",
                ],
            ),
        ],
    ),
    Deck(
        module=2,
        slug="lesson-2-1-fix-bad-email",
        title="Lesson 2.1 — Fixing a Bad Email",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 1",
                headline="Fixing a Bad Email",
                subhead="Rough draft in. Professional email out.",
            ),
            Slide(
                kind="concept",
                eyebrow="Prompting pattern",
                headline="Audience · Structure · Outcome",
                body=[
                    "•  Audience  — who's reading?",
                    "•  Structure — how should it be organized?",
                    "•  Outcome   — what action do you want from them?",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="Your turn",
                headline="Submit. Get feedback. See examples.",
                callout="Three example rewrites are waiting after you submit:\nbulleted · action-items-first · short and direct.",
            ),
        ],
    ),
    Deck(
        module=2,
        slug="lesson-2-2-summarize-thread",
        title="Lesson 2.2 — Summarizing a Long Email Thread",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 2",
                headline="Summarizing the Chaos",
                subhead="Forty-seven replies. Ninety seconds. The whole picture.",
            ),
            Slide(
                kind="concept",
                eyebrow="Pattern",
                headline="Status · Decisions · Actions",
                body=[
                    "•  Status     — where do things stand right now?",
                    "•  Decisions  — what's been agreed?",
                    "•  Actions    — who does what, by when?",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="New term",
                headline="Context window",
                callout="The maximum amount of text AI can read at once.\nFor very long threads, summarize in chunks.",
            ),
        ],
    ),
    Deck(
        module=2,
        slug="lesson-2-3-difficult-response",
        title="Lesson 2.3 — Drafting a Difficult Response",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 3",
                headline="Drafting a Difficult Response",
                subhead="Hard emails. Right tone. Every time.",
            ),
            Slide(
                kind="concept",
                eyebrow="The recipe",
                headline="Tone calibration",
                body=[
                    "1.  State the tone you want (firm but professional).",
                    "2.  Provide the facts on your side.",
                    "3.  Acknowledge their concern — don't accept blame.",
                    "4.  Propose a path forward.",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="Security warning",
                headline="Prompt Injection",
                callout="External content can hide instructions for the AI.\nAlways review AI output — especially when input came from outside.",
            ),
        ],
    ),
    # ===================================================================== #
    # MODULE 3 — Reports and Documents
    # ===================================================================== #
    Deck(
        module=3,
        slug="module-3-intro",
        title="Module 3 Intro — Reports and Documents",
        slides=[
            Slide(
                kind="title",
                eyebrow="Module 3",
                headline="Reports and Documents",
                subhead="First drafts in minutes. Polished work in an hour.",
            ),
            Slide(
                kind="concept",
                eyebrow="The shift",
                headline="Writer → Reviewer",
                body=[
                    "•  Write a report from scratch",
                    "•  Polish a rough draft",
                    "•  Build reusable templates",
                ],
            ),
        ],
    ),
    Deck(
        module=3,
        slug="lesson-3-1-report-from-scratch",
        title="Lesson 3.1 — Writing a Report from Scratch",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 1",
                headline="From Notes to Report",
                subhead="Bullet points in. Polished report out.",
            ),
            Slide(
                kind="concept",
                eyebrow="Standard structure",
                headline="Executive summary first",
                body=[
                    "Top 3-4 sentences = the bottom line.",
                    "Then: Schedule · Budget · Safety · Completed · Upcoming · Issues.",
                    "Tell AI: \"keep specific numbers from the notes.\"",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="Always check",
                headline="Hallucinations",
                callout="Anything in the report that wasn't in the notes is suspect.\nRead the draft against your source.",
            ),
        ],
    ),
    Deck(
        module=3,
        slug="lesson-3-2-rewriting-editing",
        title="Lesson 3.2 — Rewriting and Editing",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 2",
                headline="Summarize the Quincy Reports",
                subhead="Real Cashman daily reports. Nine weeks. One PDF.",
            ),
            Slide(
                kind="concept",
                eyebrow="The task",
                headline="Executive summary + thematic insights",
                body=[
                    "•  PDF spans Dec 30, 2006 → Mar 5, 2007 (~9 weeks).",
                    "•  Word files extracted, combined into one PDF for AI.",
                    "•  Summarize the season. Group insights by theme.",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="The discipline",
                headline="Source-grounding",
                callout="Daily reports are full of specific facts.\nVerify every name, number, and date AI cites — back to the source.",
            ),
        ],
    ),
    Deck(
        module=3,
        slug="lesson-3-3-templates",
        title="Lesson 3.3 — Templates and Checklists",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 3",
                headline="Templates and Checklists",
                subhead="Build it once. Use it forever.",
            ),
            Slide(
                kind="concept",
                eyebrow="The pattern",
                headline="Standard work",
                body=[
                    "Describe the document. Ask AI for a template.",
                    "Use [Brackets] for placeholder fields.",
                    "Save the prompt — not just the template.",
                ],
            ),
        ],
    ),
    # ===================================================================== #
    # MODULE 4 — Spreadsheets and Data
    # ===================================================================== #
    Deck(
        module=4,
        slug="module-4-intro",
        title="Module 4 Intro — Spreadsheets and Data",
        slides=[
            Slide(
                kind="title",
                eyebrow="Module 4",
                headline="Spreadsheets and Data",
                subhead="Stop Googling formulas. Just describe what you want.",
            ),
            Slide(
                kind="concept",
                eyebrow="Three exercises",
                headline="One real $32M project.",
                body=[
                    "1.  AI-powered formulas — on a $32M schedule of values",
                    "2.  Smarter data analysis",
                    "3.  Cleaner data",
                ],
            ),
        ],
    ),
    Deck(
        module=4,
        slug="lesson-4-1-formulas",
        title="Lesson 4.1 — AI-Powered Formulas",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 1",
                headline="$32M Schedule of Values",
                subhead="Real-shaped formulas. Real verification. Downloadable.",
            ),
            Slide(
                kind="concept",
                eyebrow="The workbook",
                headline="21 line items. 7 divisions. 8 formulas.",
                body=[
                    "•  Inputs filled in: Total Contract, Previously, This Period",
                    "•  You write: Total Billed · % Complete · Remaining · Retention · Net Earned · Status",
                    "•  Plus: SUMIFS division subtotals + project total row",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="The rule",
                headline="Always test the formula.",
                callout="$32M means a wrong formula = real-money mistake.\nTest with known rows. Project % complete is weighted, never averaged.",
            ),
        ],
    ),
    Deck(
        module=4,
        slug="lesson-4-2-data-analysis",
        title="Lesson 4.2 — Data Analysis with AI",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 2",
                headline="Data Analysis with AI",
                subhead="32 hypothetical projects. One question. Your on-call analyst.",
            ),
            Slide(
                kind="concept",
                eyebrow="Three concepts",
                headline="Pivot · Variance · Correlation",
                body=[
                    "•  Pivot table — summarize by category.",
                    "•  Variance — actual minus expected.",
                    "•  Correlation ≠ Causation. Don't confuse the two.",
                ],
            ),
            Slide(
                kind="image",
                eyebrow="Copilot in Excel — live",
                headline="Cost overrun drivers, in seconds.",
                image="module-4/copilot-excel-cost-overrun-drivers.png",
                caption=(
                    "Copilot in Excel built the Summary by Project Type, the "
                    "Cost Overrun Drivers table, and the risk profile column "
                    "from one prompt. Read it — but always verify the "
                    "numbers tie back to the source rows."
                ),
            ),
            Slide(
                kind="callout",
                eyebrow="The judgment call",
                headline="AI suggests. You decide.",
                callout="32 projects across 4 categories isn't enough for regression.\nApply your domain expertise to AI's suggestions.",
            ),
        ],
    ),
    Deck(
        module=4,
        slug="lesson-4-3-data-cleanup",
        title="Lesson 4.3 — Data Cleanup",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 3",
                headline="Closing Out a Messy Project",
                subhead="Seven files. Three deliverables. One afternoon.",
            ),
            Slide(
                kind="concept",
                eyebrow="The IPC Lydon bundle",
                headline="What you have. What you owe.",
                body=[
                    "•  Inputs — Word logs, Excel cost + payment, txt transcripts, emails.",
                    "•  Outputs — 1-page Word summary · Excel sequence-of-events · (bonus) PowerPoint.",
                    "•  Always cite the source file for every number.",
                ],
            ),
            Slide(
                kind="image",
                eyebrow="ChatGPT — the whole bundle at once",
                headline="Drop the zip. Ask one focused question.",
                image="module-4/chatgpt-ipc-lydon-zip-prompt.png",
                caption=(
                    "Attach every file from the bundle in a single message, "
                    "then ask for a chronological event list with source-file "
                    "citations. Fastest cross-file synthesis — but the data "
                    "leaves your network, so use it on hypothetical or public "
                    "data only."
                ),
            ),
            Slide(
                kind="image",
                eyebrow="Cashman AI Portal — internal data",
                headline="Same approach. On-prem.",
                image="module-4/cashman-portal-summarize.png",
                caption=(
                    "For real Cashman files, attach to the Portal instead. "
                    "Files stay on company infrastructure. Best paired with "
                    "Microsoft Copilot inside Word/Excel for per-file "
                    "follow-ups."
                ),
            ),
            Slide(
                kind="callout",
                eyebrow="The audit trail",
                headline="AI is the typist. You are the auditor.",
                callout="Every dollar in your timeline must tie to cost-tracking.xlsx.\nDates chronological. No invented quotes.",
            ),
        ],
    ),
    # ===================================================================== #
    # MODULE 5 — Images, Video, and Media
    # ===================================================================== #
    Deck(
        module=5,
        slug="module-5-intro",
        title="Module 5 Intro — Images, Video, and Media",
        slides=[
            Slide(
                kind="title",
                eyebrow="Module 5",
                headline="Images, Video, and Media",
                subhead="AI doesn't just write — it draws, designs, listens, and can fake.",
            ),
            Slide(
                kind="concept",
                eyebrow="What you'll do",
                headline="Generate. Build. Transcribe.",
                body=[
                    "1.  Generate presentation images",
                    "2.  Build slide decks from notes",
                    "3.  Transcribe and summarize meetings",
                ],
            ),
        ],
    ),
    Deck(
        module=5,
        slug="lesson-5-1-images",
        title="Lesson 5.1 — Creating Images for Presentations",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 1",
                headline="Creating Images",
                subhead="No design skills required.",
            ),
            Slide(
                kind="concept",
                eyebrow="Image-prompt framework",
                headline="Subject · Composition · Style · Exclusions",
                body=[
                    "•  Subject     — what's in the image",
                    "•  Composition — how it's framed",
                    "•  Style       — what it looks like",
                    "•  Exclusions  — what to leave out",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="Security warning",
                headline="Deepfakes",
                callout="The same tech that makes useful images makes convincing fake ones.\nVerify any photo that supports a financial or legal claim.",
            ),
        ],
    ),
    Deck(
        module=5,
        slug="lesson-5-2-presentations",
        title="Lesson 5.2 — Presentations with AI",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 2",
                headline="Presentations with AI",
                subhead="Outline first. Detail second.",
            ),
            Slide(
                kind="concept",
                eyebrow="The two-step pattern",
                headline="Progressive elaboration",
                body=[
                    "1.  Ask for the outline first — titles + one-line purpose per slide.",
                    "2.  Approve the structure. Then ask for slide content.",
                    "Tell AI exactly who the audience is.",
                ],
            ),
        ],
    ),
    Deck(
        module=5,
        slug="lesson-5-3-video-audio",
        title="Lesson 5.3 — Video, Audio, and Transcription",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 3",
                headline="Video, Audio, Transcription",
                subhead="Meetings document themselves now.",
            ),
            Slide(
                kind="concept",
                eyebrow="What's possible",
                headline="Transcribe · Summarize · Follow up",
                body=[
                    "•  Auto-transcription with speaker diarization",
                    "•  AI-generated summary, action items, draft email",
                    "•  And: voice cloning + synthetic video are real",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="The defense",
                headline="Verify through a separate channel.",
                callout="Unusual financial or legal request by voicemail or email?\nCall them back on a phone number you already have.",
            ),
        ],
    ),
    # ===================================================================== #
    # MODULE 6 — Document Processing and Research
    # (merged the old Module 8 / Power User Tools into Module 6)
    # ===================================================================== #
    Deck(
        module=6,
        slug="module-6-intro",
        title="Module 6 Intro — Document Processing and Research",
        slides=[
            Slide(
                kind="title",
                eyebrow="Module 6",
                headline="Document Processing & Research",
                subhead="Find what's inside. Research what's outside. Build the deck.",
            ),
            Slide(
                kind="concept",
                eyebrow="Three lessons",
                headline="Search · Deep Research · Workflow",
                body=[
                    "1.  Search company documents (the Portal + RAG)",
                    "2.  Deep Research → Notebook LM → AI Slides",
                    "3.  Build your AI workflow + take the survey",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="One badge away",
                headline="Think (AI)mpossible.",
                callout="Finish Module 6, take the survey, pass the final assessment.\nYou earn the badge and the certificate.",
            ),
        ],
    ),
    Deck(
        module=6,
        slug="lesson-6-1-search",
        title="Lesson 6.1 — Searching Company Documents",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 1",
                headline="Searching Company Documents",
                subhead="Stop hunting. Start asking.",
            ),
            Slide(
                kind="concept",
                eyebrow="New terms",
                headline="Embedding · Semantic Search · RAG",
                body=[
                    "Embedding — a mathematical \"fingerprint\" of meaning.",
                    "Semantic search finds meaning, not just exact words.",
                    "RAG = Retrieval-Augmented Generation. Open-book exam.",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="Always",
                headline="Verify the citation.",
                callout="A good RAG system tells you where the answer came from.\nClick through. Confirm the source actually says that.",
            ),
        ],
    ),
    Deck(
        module=6,
        slug="lesson-6-2-deep-research",
        title="Lesson 6.2 — Deep Research End-to-End",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 2",
                headline="Deep Research End-to-End",
                subhead="Three tools. One pipeline. Under an hour.",
            ),
            Slide(
                kind="concept",
                eyebrow="The pipeline",
                headline="Research → Synthesize → Slides",
                body=[
                    "1.  ChatGPT or Claude Deep Research → cited report",
                    "2.  Google Notebook LM → briefing + audio overview",
                    "3.  Claude / GPT / Copilot → 10-slide deck",
                ],
            ),
            Slide(
                kind="concept",
                eyebrow="Prompt formula",
                headline="Decision · Audience · Sources · Format",
                body=[
                    "•  Decision the report supports",
                    "•  Audience reading it",
                    "•  Source constraints (peer-reviewed, OEM, USACE)",
                    "•  Output format (markdown, sections, citations)",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="Two non-negotiables",
                headline="Public data only · Verify before you ship.",
                callout="No client names, no internal pricing in public Deep Research.\nSpot-check the biggest 2-3 numbers and click the citations.",
            ),
        ],
    ),
    Deck(
        module=6,
        slug="lesson-6-3-ai-workflow",
        title="Lesson 6.3 — Building Your AI Workflow",
        slides=[
            Slide(
                kind="title",
                eyebrow="Lesson 3",
                headline="Building Your AI Workflow",
                subhead="Wrap up. Lock it in. Take the badge.",
            ),
            Slide(
                kind="concept",
                eyebrow="Security recap",
                headline="Five risks. Five defenses.",
                body=[
                    "•  Prompt Injection  — review AI output on outside input",
                    "•  Hallucinations    — source-grounding + verification",
                    "•  Wrong Formulas    — test with known values",
                    "•  Deepfakes         — verify through a separate channel",
                    "•  Data Leakage      — company data → Portal. Period.",
                ],
            ),
            Slide(
                kind="callout",
                eyebrow="The cheat sheet",
                headline="Print it. Tape it. Live it.",
                callout="Survey → Final Assessment → Think (AI)mpossible.\nYou did the work. You earned this.",
            ),
        ],
    ),
]


# --------------------------------------------------------------------------- #
# Build
# --------------------------------------------------------------------------- #


def build_deck(deck: Deck, output_root: Path, assets_root: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for s in deck.slides:
        if s.kind == "title":
            add_title_slide(prs, s.eyebrow, s.headline, s.subhead)
        elif s.kind == "concept":
            add_concept_slide(prs, s.eyebrow, s.headline, s.body)
        elif s.kind == "callout":
            add_callout_slide(prs, s.eyebrow, s.headline, s.callout)
        elif s.kind == "image":
            img_path = assets_root / s.image
            if not img_path.exists():
                raise FileNotFoundError(f"Image asset not found: {img_path}")
            add_image_slide(prs, s.eyebrow, s.headline, img_path, s.caption)
        else:
            raise ValueError(f"Unknown slide kind: {s.kind}")

    out_dir = output_root / f"module-{deck.module}"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{deck.slug}.pptx"
    prs.save(str(out_path))
    return out_path


def main() -> None:
    repo_root = Path(__file__).resolve().parent.parent
    output_root = repo_root / "docs" / "video-scripts" / "slides" / "pptx"
    assets_root = repo_root / "docs" / "video-scripts" / "slides" / "assets"
    output_root.mkdir(parents=True, exist_ok=True)

    print(f"Generating {len(DECKS)} decks → {output_root}")
    for deck in DECKS:
        path = build_deck(deck, output_root, assets_root)
        rel = path.relative_to(repo_root)
        print(f"  ✓ {rel}  ({len(deck.slides)} slides)")
    print("Done.")


if __name__ == "__main__":
    main()
